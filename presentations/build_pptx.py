#!/usr/bin/env python3
"""Build the Patch ChatGPT workflow PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
PATCH_GREEN  = RGBColor(0x00, 0x4C, 0x05)
MID_GREEN    = RGBColor(0x00, 0x6B, 0x1A)
LIGHT_GREEN  = RGBColor(0xE8, 0xF5, 0xE9)
CODE_BG      = RGBColor(0xF0, 0xF7, 0xF0)
CODE_FG      = RGBColor(0x00, 0x60, 0x18)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1A, 0x1A, 0x1A)

# ── Fonts ──────────────────────────────────────────────────────────────────────
TF = "Roboto"       # titles
BF = "Merriweather" # body
CF = "Courier New"  # code / prompts

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE = "/home/richkirby/SpiderOak Hive/Code/GitHubProjects/Clinton-Claude/presentations/"
LOGO     = BASE + "patch_logo_cropped.jpg"   # black bars removed; 1200×456, ratio 2.63:1
ICON     = BASE + "Ptach Media Logo 01.png"   # 225×225 square
OUT      = BASE + "Patch-ChatGPT-Town-Meetings-Workflow.pptx"

# ── Presentation setup ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(BLANK)

def box(sl, l, t, w, h, fill, line_color=None):
    """Add a filled rectangle."""
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
    else:
        s.line.fill.background()
    return s

def tx(sl, text, l, t, w, h,
       font=BF, size=18, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT):
    """Single-paragraph textbox."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def mtx(sl, lines, l, t, w, h,
        font=BF, size=18, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    """Multi-line textbox.  lines may be plain strings or dicts with overrides."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            d = {}
            txt = ln
        else:
            d = ln
            txt = d.get("text", "")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = d.get("align", align)
        r = p.add_run()
        r.text = txt
        r.font.name   = d.get("font",   font)
        r.font.size   = Pt(d.get("size", size))
        r.font.bold   = d.get("bold",   bold)
        r.font.italic = d.get("italic", italic)
        r.font.color.rgb = d.get("color", color)
    return tb

def logo_small(sl):
    """Patch wordmark — bottom-right corner. Width drives height at natural 2.63:1 ratio."""
    sl.shapes.add_picture(LOGO, Inches(11.43), Inches(6.72), Inches(1.8), Inches(0.684))

def header(sl, title):
    """Green header bar + white title."""
    box(sl, 0, 0, 13.33, 1.35, PATCH_GREEN)
    tx(sl, title, 0.35, 0.12, 12.6, 1.1, font=TF, size=30, bold=True, color=WHITE)
    # thin rule at foot
    box(sl, 0, 7.42, 13.33, 0.08, PATCH_GREEN)

def step_badge(sl, num, l, t, sz=0.55):
    """Numbered circle-ish badge."""
    box(sl, l, t, sz, sz, PATCH_GREEN)
    tx(sl, str(num), l, t, sz, sz, font=TF, size=22, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, PATCH_GREEN)
box(s, 0, 5.7, 13.33, 0.07, WHITE)

tx(s,
   "How Patch Journalists Can Use ChatGPT\nTo Turn Town Meetings Into News Stories",
   0.8, 1.1, 11.73, 3.4,
   font=TF, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tx(s,
   "A practical step-by-step workflow for covering local government meetings\n"
   "using ChatGPT, transcription tools, and Associated Press-style prompts.",
   1.2, 4.1, 10.93, 1.5,
   font=BF, size=18, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

s.shapes.add_picture(LOGO, Inches(4.665), Inches(5.65), Inches(4.0), Inches(1.52))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What Kind of Meetings?
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "What Kind of Meetings Are We Talking About?")
logo_small(s)

tx(s,
   "The resident volunteers on these boards and commissions meet regularly "
   "to hash out matters that often keenly affect your readership:",
   0.5, 1.48, 12.33, 0.85, size=18)

boards = [
    "Boards of Selectmen / Town Council",
    "Planning & Zoning Commission",
    "Board of Finance",
    "Board of Education",
    "Inland Wetlands Commission",
    "Zoning Board of Appeals",
    "…and many other local committees and commissions",
]
mtx(s, [{"text": f"  •  {b}", "size": 17} for b in boards],
    0.9, 2.4, 11.8, 4.8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why Write About Municipal Meetings?
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Why Write About Municipal Meetings?")
logo_small(s)

# Left panel
box(s, 0.4, 1.45, 5.95, 5.75, LIGHT_GREEN)
tx(s, "The Challenge", 0.6, 1.55, 5.55, 0.55, font=TF, size=20, bold=True, color=PATCH_GREEN)
tx(s,
   "Municipal meetings have not exactly proven to be pageview "
   "magnets — they require time, expertise, and consistent "
   "attendance to cover well.",
   0.6, 2.15, 5.55, 1.8, size=16)

# Right panel
box(s, 6.98, 1.45, 6.0, 5.75, PATCH_GREEN)
tx(s, "The Opportunity", 7.18, 1.55, 5.6, 0.55, font=TF, size=20, bold=True, color=WHITE)

opps = [
    "Establish Patch as the local government authority",
    "SEO → GEO: LLMs reward perceived authority with superior link placement",
    'Tag stories as "Patch Exclusives"',
    "Cover news your competition is overlooking",
    "AI makes deep local reporting possible for busy editors",
]
mtx(s, [{"text": f"  ✓  {o}", "size": 15, "color": WHITE} for o in opps],
    7.18, 2.2, 5.6, 4.8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SEO → GEO
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "SEO Is Becoming GEO")
logo_small(s)

tx(s,
   "Generative Engine Optimization (GEO) is the new frontier for local news visibility.",
   0.5, 1.5, 12.33, 0.65, font=TF, size=20, bold=True, color=PATCH_GREEN)

tx(s,
   "As AI search (Google AI Overviews, ChatGPT, Perplexity) replaces traditional "
   "results, LLMs reward the perceived authority of a news outlet with superior — "
   "if not exclusive — link placement in their responses.",
   0.5, 2.25, 12.33, 1.5, size=18)

tx(s, "To earn that authority:", 0.5, 3.9, 12.33, 0.5, size=18, bold=True)

pts = [
    "Provide important local news that competitors overlook",
    "Cover local government meetings consistently",
    "Publish on a regular, predictable basis",
]
mtx(s, [{"text": f"  •  {p}", "size": 17} for p in pts],
    0.9, 4.45, 12.0, 1.6)

tx(s,
   "AI makes all of this achievable — even for editors managing multiple Patches.",
   0.5, 6.2, 12.33, 0.65, size=18, italic=True, color=PATCH_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Workflow 1 Overview
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, LIGHT_GREEN)
box(s, 0, 0, 13.33, 0.1, PATCH_GREEN)
box(s, 0, 7.4, 13.33, 0.1, PATCH_GREEN)
logo_small(s)

tx(s, "WORKFLOW 1", 0.5, 0.7, 12.33, 0.7,
   font=TF, size=22, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)
tx(s, "Meeting Minutes  →  News Article", 0.5, 1.35, 12.33, 1.4,
   font=TF, size=50, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)

steps = ["Get\nMinutes", "Open\nChatGPT", "Paste\nPrompt", "Attach\nMinutes", "Review\nDraft", "Publish"]
x = 0.55
W_BOX = 1.88
for i, step in enumerate(steps, 1):
    box(s, x, 3.8, W_BOX, 2.1, PATCH_GREEN)
    tx(s, str(i), x, 3.82, W_BOX, 0.78,
       font=TF, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mtx(s, [{"text": step, "align": PP_ALIGN.CENTER, "size": 14, "color": WHITE, "font": BF}],
        x + 0.08, 4.62, W_BOX - 0.16, 1.2)
    if i < 6:
        tx(s, "›", x + W_BOX + 0.02, 4.35, 0.28, 0.7,
           font=TF, size=26, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)
    x += W_BOX + 0.28


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — W1 Steps 1-3
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 1  |  Steps 1–3")
logo_small(s)

rows = [
    (1, 1.48, "Get the Meeting Minutes",
     "Download or copy the official meeting minutes from the town website. "
     "Many towns offer free email subscriptions to receive minutes and agendas "
     "(usually as PDFs or MS Word documents)."),
    (2, 3.1,  "Open ChatGPT",
     "Start a new chat in ChatGPT."),
    (3, 4.3,  "Paste the Prompt",
     "Copy and paste the AP-style prompt into ChatGPT. "
     "See the next slide for the complete prompt text."),
]
for num, top, title, body in rows:
    step_badge(s, num, 0.4, top)
    tx(s, title, 1.1, top + 0.02, 11.8, 0.52,
       font=TF, size=21, bold=True, color=PATCH_GREEN)
    tx(s, body, 1.1, top + 0.57, 11.8, 1.0, size=17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — W1 Prompt
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 1  |  The Prompt  (copy & paste this)")
logo_small(s)

box(s, 0.4, 1.42, 12.53, 5.78, CODE_BG)

prompt1 = [
    "Rewrite these [DATE] meeting minutes of the [TOWN], [STATE]",
    "[MUNICIPAL BODY] as a news article in the neutral style of an",
    "Associated Press journalist.",
    "",
    "Requirements:",
    "    Lead with the most newsworthy actions or votes.",
    "    Attribute comments clearly.",
    "    Summarize discussion concisely.",
    "    Include vote counts when available.",
    "    Do not editorialize or speculate.",
    "    Write for a hyperlocal readership in [TOWN], [STATE].",
    "",
    "Then generate:",
    "    An SEO-friendly headline (≤109 characters)",
    "    A meta description (≤136 characters)",
    "",
    "Meeting minutes:",
    "    [PASTE TEXT HERE, or attach the file to this prompt]",
]
mtx(s, [{"text": ln, "font": CF, "size": 14, "color": CODE_FG} for ln in prompt1],
    0.6, 1.55, 12.1, 5.55, wrap=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — W1 Steps 4-6
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 1  |  Steps 4–6")
logo_small(s)

# Step 4
step_badge(s, 4, 0.4, 1.48)
tx(s, "Attach or Paste the Minutes", 1.1, 1.5, 11.8, 0.52, font=TF, size=21, bold=True, color=PATCH_GREEN)
tx(s, "Paste the full text directly into the prompt, or upload the minutes file as an attachment.",
   1.1, 2.05, 11.8, 0.55, size=17)

# Step 5
step_badge(s, 5, 0.4, 2.78)
tx(s, "Review the Draft Carefully", 1.1, 2.8, 11.8, 0.52, font=TF, size=21, bold=True, color=PATCH_GREEN)
tx(s, "Paste ChatGPT's work into the Patch CMS. Before publishing:", 1.1, 3.35, 11.8, 0.45, size=17)

checks = ["Verify vote counts", "Confirm names and titles",
          "Check quotes against the original minutes",
          "Make sure the lede reflects the most important action",
          "Remove anything that sounds promotional or speculative"]
mtx(s, [{"text": f"    ✓  {c}", "size": 16, "color": PATCH_GREEN} for c in checks],
    1.1, 3.85, 11.8, 2.1)

# Step 6
step_badge(s, 6, 0.4, 6.05)
tx(s, "Publish", 1.1, 6.07, 11.8, 0.5, font=TF, size=21, bold=True, color=PATCH_GREEN)
tx(s,
   "Tweak the hed, lede and meta description. Add pull-quotes to make the story yours, "
   "and 2–3 recirc links to previous meeting coverage.",
   1.1, 6.6, 11.8, 0.75, size=16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Agenda Variation
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Bonus: The Agenda Variation  —  The 2-for-1 Play")
logo_small(s)

tx(s,
   "You can often squeeze a forward-looking story from a committee's published "
   "agenda — before the meeting even happens.",
   0.5, 1.5, 12.33, 0.8, size=18)

box(s, 0.4, 2.45, 12.53, 3.05, CODE_BG)
tx(s, "How to modify the prompt:", 0.6, 2.52, 12.0, 0.5,
   font=TF, size=17, bold=True, color=PATCH_GREEN)

mods = [
    "1.  Remove all \"Requirements\" except:",
    "        Write for a hyperlocal readership in [TOWN], [STATE].",
    "",
    "2.  Change:    these [DATE] meeting minutes",
    "    To:              this [DATE] meeting agenda",
]
mtx(s, [{"text": ln, "font": CF, "size": 14, "color": CODE_FG} for ln in mods],
    0.65, 3.05, 12.0, 2.3, wrap=False)

tx(s,
   "The AI will create a forward-looking meeting preview. "
   "Consider scheduling it to publish the day before or day of the meeting — "
   "with an appropriately urgent headline.",
   0.5, 5.6, 12.33, 1.25, size=17, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Workflow 2 Overview
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, LIGHT_GREEN)
box(s, 0, 0, 13.33, 0.1, PATCH_GREEN)
box(s, 0, 7.4, 13.33, 0.1, PATCH_GREEN)
logo_small(s)

tx(s, "WORKFLOW 2", 0.5, 0.7, 12.33, 0.7,
   font=TF, size=22, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)
tx(s, "Meeting Transcript  →  News Article", 0.5, 1.35, 12.33, 1.4,
   font=TF, size=46, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)
tx(s,
   "Use when you have audio or video. Results are invariably superior to working from minutes alone.",
   1.2, 2.82, 10.93, 0.75,
   font=BF, size=17, italic=True, color=DARK, align=PP_ALIGN.CENTER)

steps2 = ["Get\nRecording", "Create\nTranscript", "Export\nTranscript",
          "Open\nChatGPT", "Paste\nPrompt", "Fact-\nCheck", "Add\nContext"]
x = 0.22
W_BOX2 = 1.74
for i, step in enumerate(steps2, 1):
    box(s, x, 3.75, W_BOX2, 2.15, PATCH_GREEN)
    tx(s, str(i), x, 3.77, W_BOX2, 0.78,
       font=TF, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mtx(s, [{"text": step, "align": PP_ALIGN.CENTER, "size": 13, "color": WHITE, "font": BF}],
        x + 0.07, 4.57, W_BOX2 - 0.14, 1.25)
    if i < 7:
        tx(s, "›", x + W_BOX2 + 0.02, 4.3, 0.24, 0.7,
           font=TF, size=22, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)
    x += W_BOX2 + 0.24


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — W2 Steps 1-2
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 2  |  Steps 1–2")
logo_small(s)

step_badge(s, 1, 0.4, 1.48)
tx(s, "Obtain the Meeting Recording", 1.1, 1.5, 11.8, 0.52,
   font=TF, size=21, bold=True, color=PATCH_GREEN)
tx(s,
   "Look for: livestream recordings, YouTube uploads, Zoom recordings, MP3/MP4 files. "
   "Your town may host these on its website or link to third-party providers.",
   1.1, 2.06, 11.8, 0.8, size=17)

step_badge(s, 2, 0.4, 3.05)
tx(s, "Create a Transcript", 1.1, 3.07, 11.8, 0.52,
   font=TF, size=21, bold=True, color=PATCH_GREEN)

# Two option panels
box(s, 1.1, 3.68, 5.4, 3.3, CODE_BG)
tx(s, "Option A: Upload Directly", 1.25, 3.76, 5.1, 0.52,
   font=TF, size=17, bold=True, color=PATCH_GREEN)
tx(s,
   "If you have a standalone audio/video file:\n"
   "  1. Upload to Otter.ai, TurboScribe, or another service\n"
   "  2. Wait for transcription to complete",
   1.25, 4.32, 5.1, 2.5, size=15)

box(s, 6.9, 3.68, 6.0, 3.3, CODE_BG)
tx(s, "Option B: Record Playback Into Otter.ai", 7.05, 3.76, 5.7, 0.52,
   font=TF, size=17, bold=True, color=PATCH_GREEN)
tx(s,
   "If the meeting is only on your town's website:\n"
   "  1. Open the meeting on your laptop\n"
   "  2. Open Otter.ai\n"
   "  3. Play the audio aloud\n"
   "  4. Let Otter.ai transcribe in real time",
   7.05, 4.32, 5.7, 2.5, size=15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — W2 Steps 3-5
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 2  |  Steps 3–5")
logo_small(s)

rows2 = [
    (3, 1.48, "Export the Transcript",
     "Export as plain text (.txt) — this format works best with ChatGPT."),
    (4, 2.88, "Open ChatGPT",
     "Start a new chat."),
    (5, 4.05, "Upload the Transcript & Paste the Prompt",
     "Attach the transcript file to the prompt, then paste the AP-style prompt "
     "(full text on next slide). ChatGPT will process both together."),
]
for num, top, title, body in rows2:
    step_badge(s, num, 0.4, top)
    tx(s, title, 1.1, top + 0.02, 11.8, 0.52,
       font=TF, size=21, bold=True, color=PATCH_GREEN)
    tx(s, body, 1.1, top + 0.57, 11.8, 1.0, size=17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — W2 Prompt
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 2  |  The Prompt  (copy & paste this)")
logo_small(s)

box(s, 0.4, 1.42, 12.53, 5.78, CODE_BG)

prompt2 = [
    "Rewrite this transcript of the [DATE] meeting of the [TOWN],",
    "[STATE] [MUNICIPAL BODY] as a news article in the neutral style of",
    "an Associated Press journalist.",
    "",
    "Requirements:",
    "    Identify speakers by name and title when possible.",
    "    Remove filler, repetition and off-topic remarks.",
    "    Highlight decisions, disagreements and next steps.",
    "    Preserve exact meaning of quotes.",
    "    Do not invent attribution.",
    "    Write for a hyperlocal readership in [TOWN], [STATE].",
    "",
    "Then generate:",
    "    1) An SEO-friendly headline (≤109 characters)",
    "    2) A meta description (≤136 characters)",
    "",
    "Attach the transcript file to this prompt.",
]
mtx(s, [{"text": ln, "font": CF, "size": 14, "color": CODE_FG} for ln in prompt2],
    0.6, 1.55, 12.1, 5.55, wrap=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — W2 Steps 6-7
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Workflow 2  |  Steps 6–7")
logo_small(s)

step_badge(s, 6, 0.4, 1.48)
tx(s, "Fact-Check the Draft", 1.1, 1.5, 11.8, 0.52,
   font=TF, size=21, bold=True, color=PATCH_GREEN)
tx(s, "Paste ChatGPT's work into the Patch CMS. Always verify:",
   1.1, 2.05, 11.8, 0.45, size=17)

checks2 = ["Names and titles", "Vote counts", "Quotes",
           "Timeline of decisions", "Motions and outcomes"]
mtx(s, [{"text": f"    ✓  {c}", "size": 16, "color": PATCH_GREEN} for c in checks2],
    1.1, 2.55, 6.0, 2.1)

tx(s, "AI can miss context or confuse speakers during complicated discussions.",
   1.1, 4.72, 11.8, 0.5, size=15, italic=True)

step_badge(s, 7, 0.4, 5.3)
tx(s, "Add Context Before Publishing", 1.1, 5.32, 11.8, 0.52,
   font=TF, size=21, bold=True, color=PATCH_GREEN)
ctx = ["Prior meeting history", "Budget context",
       "Links to applications or agenda packets",
       "Previous controversies",
       "Public reaction from 'public comment' portions"]
mtx(s, [{"text": f"  •  {c}", "size": 16} for c in ctx],
    1.1, 5.88, 11.8, 1.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Tips for Better Results
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Tips for Better Results")
logo_small(s)

# Left panel — follow-up rewrites
box(s, 0.4, 1.45, 5.9, 5.75, LIGHT_GREEN)
tx(s, "Ask for Follow-Up Rewrites", 0.6, 1.53, 5.5, 0.55,
   font=TF, size=19, bold=True, color=PATCH_GREEN)

rewrites = [
    '"Rewrite with a stronger lede"',
    '"Shorten this by 200 words"',
    '"Focus more on the budget vote"',
    '"Organize chronologically"',
    '"Highlight the major disagreements"',
]
mtx(s, [{"text": f"  •  {r}", "font": CF, "size": 15, "color": CODE_FG} for r in rewrites],
    0.6, 2.15, 5.5, 2.8)

# Right panel — upload docs
box(s, 7.03, 1.45, 5.9, 5.75, PATCH_GREEN)
tx(s, "Upload Supporting Documents", 7.23, 1.53, 5.5, 0.55,
   font=TF, size=19, bold=True, color=WHITE)
tx(s,
   "Context and accuracy improve significantly when you also attach:",
   7.23, 2.15, 5.5, 0.75, size=16, color=WHITE)

docs = ["Agendas", "Meeting packets", "Presentations",
        "Proposed budgets", "Prior meeting minutes"]
mtx(s, [{"text": f"  •  {d}", "size": 16, "color": WHITE} for d in docs],
    7.23, 2.98, 5.5, 2.2)

tx(s, "Identify each document by name inside your prompt.",
   7.23, 5.35, 5.5, 0.7, size=14, italic=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Editorial Reminder
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, PATCH_GREEN)

tx(s, "Important Editorial Reminder", 0.5, 0.45, 12.33, 0.85,
   font=TF, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tx(s,
   "ChatGPT is a reporting assistant — not a replacement for editorial judgment.",
   0.9, 1.38, 11.53, 0.75,
   font=BF, size=22, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 1.5, 2.3, 10.33, 3.85, MID_GREEN)

tx(s, "Always:", 2.0, 2.45, 10.0, 0.52, font=TF, size=20, bold=True, color=WHITE)

always = ["Verify facts", "Review quotes", "Confirm names and titles",
          "Read the original source material",
          "Apply newsroom standards before publication"]
mtx(s, [{"text": f"  ✓  {a}", "size": 18, "color": WHITE} for a in always],
    2.0, 3.02, 9.5, 3.0)

tx(s,
   "The final responsibility for accuracy remains with the journalist.",
   0.9, 5.92, 11.53, 0.65,
   font=BF, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

s.shapes.add_picture(LOGO, Inches(5.415), Inches(6.47), Inches(2.5), Inches(0.95))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Closing
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, LIGHT_GREEN)
box(s, 0, 0, 13.33, 0.1, PATCH_GREEN)
box(s, 0, 7.4, 13.33, 0.1, PATCH_GREEN)

tx(s, "Now Go Cover Those Meetings.", 0.5, 1.6, 12.33, 1.6,
   font=TF, size=52, bold=True, color=PATCH_GREEN, align=PP_ALIGN.CENTER)

tx(s,
   "Local government journalism matters — and with these tools,\n"
   "you can do it faster, smarter and more consistently.",
   1.5, 3.45, 10.33, 1.4,
   font=BF, size=20, italic=True, color=DARK, align=PP_ALIGN.CENTER)

s.shapes.add_picture(LOGO, Inches(4.665), Inches(5.0), Inches(4.0), Inches(1.52))
s.shapes.add_picture(ICON, Inches(6.265), Inches(6.6),  Inches(0.8), Inches(0.8))


# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
